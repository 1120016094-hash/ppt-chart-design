#!/usr/bin/env python3
"""Place rendered chart image(s) onto a 16:9 PPTX slide (high-fidelity image mode).

The chart is inserted as a picture - pixel-faithful to the design, not an editable
native PPT chart. Background can be set to match the style (e.g. full-bleed colored
page for Style 5).

Usage:
    python build_pptx.py --image chart.png --out deck.pptx
    python build_pptx.py --image chart.png --out deck.pptx --bg "#006CC0" --mode fullbleed
"""
import argparse
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)


def _rgb(hex_str):
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def add_chart_slide(prs, image_path, bg=None, mode="framed"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    if bg:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(bg)
    from PIL import Image
    iw, ih = Image.open(image_path).size
    ar = iw / ih
    if mode == "fullbleed":
        # cover the slide, center-crop by overflow
        if ar > SLIDE_W / SLIDE_H:
            h = SLIDE_H; w = Emu(int(h * ar))
        else:
            w = SLIDE_W; h = Emu(int(w / ar))
        left = Emu(int((SLIDE_W - w) / 2)); top = Emu(int((SLIDE_H - h) / 2))
    else:  # framed: fit with margin
        margin = Inches(0.6)
        avail_w = SLIDE_W - 2 * margin; avail_h = SLIDE_H - 2 * margin
        if ar > avail_w / avail_h:
            w = avail_w; h = Emu(int(w / ar))
        else:
            h = avail_h; w = Emu(int(h * ar))
        left = Emu(int((SLIDE_W - w) / 2)); top = Emu(int((SLIDE_H - h) / 2))
    slide.shapes.add_picture(image_path, left, top, width=w, height=h)
    return slide


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--image", required=True, action="append",
                    help="chart image path (repeat for multiple slides)")
    ap.add_argument("--out", required=True)
    ap.add_argument("--bg", default=None, help="slide background hex, e.g. #006CC0")
    ap.add_argument("--mode", default="framed", choices=["framed", "fullbleed"])
    a = ap.parse_args()
    prs = Presentation()
    prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    for img in a.image:
        add_chart_slide(prs, img, bg=a.bg, mode=a.mode)
    prs.save(a.out)
    print("wrote", a.out, "with", len(a.image), "slide(s)")


if __name__ == "__main__":
    main()
